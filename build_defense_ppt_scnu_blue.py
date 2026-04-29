# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


BASE = Path(r"D:\GNN_Thesis")
OUT = BASE / "基于时空图神经网络的多行为序列推荐模型答辩ppt_蓝色版_校徽版_v6.pptx"
LOGO = BASE / "华师标识.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(241, 246, 252)
NAVY = RGBColor(18, 59, 110)
BLUE = RGBColor(39, 96, 167)
LIGHT = RGBColor(223, 235, 248)
LINE = RGBColor(188, 209, 232)
TEXT = RGBColor(47, 63, 84)
MUTED = RGBColor(102, 121, 145)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, idx: int):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height)
    left.fill.solid()
    left.fill.fore_color.rgb = NAVY
    left.line.fill.background()

    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.72), Inches(0.6), Inches(0.14), Inches(5.9))
    deco.fill.solid()
    deco.fill.fore_color.rgb = LIGHT
    deco.line.fill.background()

    if LOGO.exists():
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(0.20), Inches(3.55), Inches(0.74))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LINE
        slide.shapes.add_picture(str(LOGO), Inches(8.65), Inches(0.24), Inches(3.32), Inches(0.64))

    page = slide.shapes.add_textbox(Inches(11.85), Inches(6.90), Inches(0.75), Inches(0.22))
    p = page.text_frame.paragraphs[0]
    p.text = f"{idx}/10"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, text: str):
    tb = slide.shapes.add_textbox(Inches(0.78), Inches(0.34), Inches(7.2), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(0.92), Inches(1.0), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_text(slide, text: str, x, y, w, h, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb


def add_card(slide, title: str, body: str, x, y, w, h, fill=WHITE, title_size=15, body_size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    add_text(slide, title, x + 0.16, y + 0.10, w - 0.3, 0.24, title_size, BLUE, True)
    add_text(slide, body, x + 0.16, y + 0.38, w - 0.3, h - 0.46, body_size, TEXT)


def add_bullets(slide, items, x, y, w, h, size=18):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
    return tb


def add_image(slide, name: str, x, y, w, h):
    path = BASE / name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def slide_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 1)
    add_text(s, "基于时空图神经网络的\n多行为序列推荐模型", 1.55, 1.35, 10.1, 1.45, 29, NAVY, True, PP_ALIGN.CENTER)
    add_text(s, "本科毕业论文答辩", 4.1, 3.08, 5.0, 0.36, 20, BLUE, True, PP_ALIGN.CENTER)
    add_text(s, "作者：黎鋆", 4.35, 4.10, 1.8, 0.30, 18, TEXT, True, PP_ALIGN.CENTER)
    add_text(s, "指导老师：董梁", 6.35, 4.10, 2.4, 0.30, 18, TEXT, True, PP_ALIGN.CENTER)
    add_text(s, "华南师范大学", 4.35, 4.70, 4.6, 0.30, 18, BLUE, True, PP_ALIGN.CENTER)


def slide_agenda():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 2)
    add_title(s, "目录")
    items = [
        ("1", "研究背景与技术路线"),
        ("2", "理论基础与数据准备"),
        ("3", "时空异构图构建"),
        ("4", "模型设计与训练策略"),
        ("5", "系统实现与软件封装"),
        ("6", "实验结果与分析"),
        ("7", "创新点与总结"),
    ]
    start_y = 1.58
    gap = 0.60
    for idx, (num, label) in enumerate(items):
        y = start_y + idx * gap
        row = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), Inches(y - 0.05), Inches(8.55), Inches(0.48))
        row.fill.solid()
        row.fill.fore_color.rgb = WHITE
        row.line.color.rgb = LINE
        row.line.width = Pt(1)
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.65), Inches(y), Inches(0.36), Inches(0.36))
        sq.fill.solid()
        sq.fill.fore_color.rgb = BLUE
        sq.line.fill.background()
        add_text(s, num, 2.65, y + 0.01, 0.36, 0.28, 17, WHITE, True, PP_ALIGN.CENTER)
        add_text(s, label, 3.22, y - 0.02, 6.75, 0.34, 23, TEXT, True, PP_ALIGN.LEFT)


def slide_background():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 3)
    add_title(s, "研究背景与技术路线")
    add_card(s, "问题来源", "电商用户行为包含点击、收藏、加购、购买等多个阶段，单一交互信号难以完整表达兴趣形成与增强过程。", 0.82, 1.40, 5.35, 1.2, LIGHT)
    add_card(s, "研究切入点", "序列模型擅长时间顺序建模，图模型擅长结构关系建模，本文希望把两者统一到同一时空图框架中。", 0.82, 2.90, 5.35, 1.2, WHITE)
    add_card(s, "技术路线", "数据清洗与字段统一 → 时空异构图构建 → 时空图神经网络训练 → 桌面系统封装 → 实验验证与诊断分析", 0.82, 4.40, 5.35, 1.05, LIGHT)
    add_image(s, "fig03_behavior_distribution.png", 6.72, 1.55, 5.15, 2.55)
    add_text(s, "多行为类型分布", 8.35, 4.18, 2.1, 0.2, 12, BLUE, True, PP_ALIGN.CENTER)
    add_card(s, "研究目标", "在统一建模框架下，同时刻画结构依赖、时间动态和行为价值差异，使推荐结果既能更准确排序，也能支持交互后的实时刷新。", 6.55, 4.48, 5.55, 0.95, WHITE)


def slide_theory_data():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 4)
    add_title(s, "理论基础与数据准备")
    add_card(s, "Top-K任务", "R_u^K = TopK_{i∈I} f(u,i)\n推荐目标是从候选集合中完成排序，而不是逐对做二分类判断。", 0.82, 1.42, 3.55, 1.18, LIGHT)
    add_card(s, "统一字段", "e=(u,i,b,t)\n将用户、商品、行为类型与时间戳统一表示，作为后续序列建模、图构建和系统接入的基础协议。", 4.55, 1.42, 3.55, 1.18, WHITE)
    add_card(s, "评价指标", "采用 HR、NDCG、MRR 衡量命中率与前列排序质量，并区分过滤已见与不过滤已见两类评估设置。", 8.28, 1.42, 3.55, 1.18, LIGHT)
    add_image(s, "fig04_user_activity_distribution.png", 0.96, 3.08, 5.02, 1.95)
    add_image(s, "fig05_item_popularity_distribution.png", 6.56, 3.08, 5.02, 1.95)
    add_text(s, "用户活跃度分布", 2.35, 5.10, 2.0, 0.2, 12, BLUE, True, PP_ALIGN.CENTER)
    add_text(s, "商品流行度分布", 7.95, 5.10, 2.0, 0.2, 12, BLUE, True, PP_ALIGN.CENTER)
    add_card(s, "数据特征", "数据呈现明显长尾分布，点击远多于购买，因此后续模型需要同时处理高频弱意图行为与低频高意图行为的不均衡问题。", 0.82, 5.48, 11.55, 0.82, WHITE)


def slide_graph():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 5)
    add_title(s, "时空异构图构建")
    add_card(s, "图定义", "G=(V,E,φ,ψ)\n将用户、商品及其交互关系组织成带类型的时空异构图，并保留行为语义与时间上下文。", 0.82, 1.42, 3.8, 1.18, LIGHT)
    add_card(s, "时间窗口", "W_k=[t_k, t_k+Δ)\n按时间窗口构建局部子图，降低长跨度噪声对当前兴趣状态的干扰。", 4.82, 1.42, 3.45, 1.18, WHITE)
    add_card(s, "行为强度", "ω(u,i,t)=w_b·exp(-λΔt)\n同时量化行为意图强弱与时间新鲜性，用于边权与样本权重计算。", 8.52, 1.42, 3.55, 1.18, LIGHT)
    add_image(s, "fig06_behavior_transition_heatmap.png", 1.05, 2.98, 5.2, 2.28)
    add_text(s, "多行为转移矩阵热力图", 2.12, 5.37, 3.0, 0.2, 12, BLUE, True, PP_ALIGN.CENTER)
    add_card(s, "构建结果", "原始行为日志经过字段规范、时间修正、行为转移统计和动态图快照组织后，被转化为可直接进入模型训练的时空图输入。", 6.62, 3.08, 5.48, 1.42, WHITE)


def slide_model():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 6)
    add_title(s, "模型设计与训练策略")
    add_image(s, "fig01_stgnn_architecture.png", 0.92, 1.42, 5.7, 3.12)
    add_card(s, "空间建模", "通过用户—商品图和商品转移图进行双图传播，学习用户与商品之间的高阶协同关系。", 6.88, 1.42, 5.15, 0.92, LIGHT)
    add_card(s, "时间建模", "将商品表示、行为嵌入、位置编码和时间间隔编码共同输入 GRU，刻画用户兴趣随行为序列变化的过程。", 6.88, 2.58, 5.15, 1.02, WHITE)
    add_card(s, "时空融合", "z_u = g_u⊙h_u^sp + (1-g_u)⊙h_u^tm\n门控机制根据当前上下文自适应决定更依赖空间结构还是时间动态。", 6.88, 3.82, 5.15, 1.02, LIGHT)
    add_card(s, "训练策略", "预测层采用内积打分，优化目标使用 BPR 排序损失，并结合随机负采样、难负样本补充和在线局部更新机制。", 0.82, 5.18, 11.2, 0.88, WHITE)


def slide_system():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 7)
    add_title(s, "系统实现与软件封装")
    add_card(s, "系统架构", "系统由数据层、模型层、服务层和界面层组成，整体遵循“界面触发—服务调度—引擎执行—结果回传”的流程。", 0.82, 1.42, 5.3, 1.02, LIGHT)
    add_card(s, "推荐引擎", "引擎统一处理训练、推理、增量更新和历史淘汰策略，离线路径负责训练，在线路径负责交互后的局部刷新。", 6.38, 1.42, 5.7, 1.02, WHITE)
    add_image(s, "双工作区_Demo.png", 0.92, 2.88, 5.7, 2.45)
    add_image(s, "推荐刷新-后.png", 6.88, 2.88, 5.25, 2.45)
    add_text(s, "双工作区界面", 2.82, 5.42, 1.9, 0.18, 12, BLUE, True, PP_ALIGN.CENTER)
    add_text(s, "实时推荐刷新", 8.78, 5.42, 1.9, 0.18, 12, BLUE, True, PP_ALIGN.CENTER)


def slide_experiment():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 8)
    add_title(s, "实验结果与分析")
    add_image(s, "fig08_main_benchmark_metrics.png", 0.92, 1.45, 5.3, 2.2)
    add_image(s, "fig09_ablation_metrics.png", 6.88, 1.45, 5.0, 2.2)
    add_text(s, "主实验结果", 2.55, 3.74, 1.8, 0.18, 12, BLUE, True, PP_ALIGN.CENTER)
    add_text(s, "消融实验结果", 8.45, 3.74, 1.8, 0.18, 12, BLUE, True, PP_ALIGN.CENTER)
    add_card(s, "实验结论", "时空联合建模在结构关系、兴趣变化和行为价值表达之间取得更平衡的效果。空间、时间和行为强度三个模块联合使用时性能最稳定。", 0.82, 4.12, 5.38, 1.18, LIGHT)
    add_card(s, "案例分析", "结合批量评分、专项诊断和用户行为可视化，可以进一步解释推荐变化原因、失败样本归因以及系统在真实交互中的响应表现。", 6.48, 4.12, 5.62, 1.18, WHITE)


def slide_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 9)
    add_title(s, "创新点与总结")
    add_card(s, "方法层", "构建了面向多行为推荐的时空统一建模框架，将结构关系、时间动态和行为强度放入同一模型联合学习。", 0.82, 1.48, 3.72, 1.45, LIGHT)
    add_card(s, "系统层", "完成桌面应用封装，实现双工作区、实时推荐刷新、自定义数据库接入、批量评分与行为可视化。", 4.82, 1.48, 3.72, 1.45, WHITE)
    add_card(s, "结果层", "通过主实验、消融实验和案例分析验证模型的有效性，并形成从数据处理到软件落地的完整研究链条。", 8.82, 1.48, 3.30, 1.45, LIGHT)
    add_bullets(s, [
        "统一建模用户—商品结构依赖、行为语义差异和时间动态变化。",
        "实现可运行、可交互、可诊断的推荐系统原型。",
        "后续可继续引入文本、图像等多模态信息，增强冷启动和长尾场景表现。"
    ], 1.05, 3.75, 10.9, 1.95, 20)


def slide_thanks():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 10)
    add_text(s, "感谢聆听", 3.15, 2.35, 7.2, 0.9, 42, NAVY, True, PP_ALIGN.CENTER)
    add_text(s, "恳请各位老师批评指正", 2.75, 3.65, 8.0, 0.4, 24, BLUE, True, PP_ALIGN.CENTER)


for fn in [
    slide_cover,
    slide_agenda,
    slide_background,
    slide_theory_data,
    slide_graph,
    slide_model,
    slide_system,
    slide_experiment,
    slide_summary,
    slide_thanks,
]:
    fn()

prs.save(OUT)
print(OUT)
