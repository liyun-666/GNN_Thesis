# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(r"D:\GNN_Thesis")
PRIMARY_OUT = BASE / "基于时空图神经网络的多行为序列推荐模型答辩ppt.pptx"
FALLBACK_OUT = BASE / "基于时空图神经网络的多行为序列推荐模型答辩ppt_10页版.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(248, 250, 252)
GREEN = RGBColor(0, 112, 84)
DARK = RGBColor(26, 43, 52)
MUTED = RGBColor(95, 108, 122)
LIGHT = RGBColor(232, 244, 240)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(208, 225, 220)
GOLD = RGBColor(196, 151, 54)


def textbox(slide, text, x, y, w, h, size=18, color=DARK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def bullets(slide, items, x, y, w, h, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    return box


def panel(slide, title_text, body, x, y, w, h, fill=WHITE, title_size=15, body_size=13):
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    textbox(slide, title_text, x + 0.14, y + 0.08, w - 0.28, 0.28, title_size, GREEN, True)
    textbox(slide, body, x + 0.14, y + 0.38, w - 0.28, h - 0.44, body_size, DARK)


def school_badge(slide, x=9.15, y=0.22, w=2.85, h=0.38):
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = RGBColor(190, 214, 208)
    shp.line.width = Pt(1)
    textbox(slide, "华南师范大学", x + 0.08, y + 0.05, w - 0.16, h - 0.08, 14, GREEN, True, PP_ALIGN.CENTER)


def background(slide, index=None):
    base = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    base.fill.solid()
    base.fill.fore_color.rgb = BG
    base.line.fill.background()

    top = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.14))
    top.fill.solid()
    top.fill.fore_color.rgb = GREEN
    top.line.fill.background()

    left = slide.shapes.add_shape(1, 0, 0, Inches(0.18), prs.slide_height)
    left.fill.solid()
    left.fill.fore_color.rgb = GREEN
    left.line.fill.background()

    deco = slide.shapes.add_shape(9, Inches(11.45), Inches(0.34), Inches(1.45), Inches(1.45))
    deco.fill.solid()
    deco.fill.fore_color.rgb = LIGHT
    deco.line.fill.background()

    school_badge(slide)

    if index is not None:
        textbox(slide, f"{index}/10", 12.0, 6.95, 0.8, 0.22, 10, MUTED, False, PP_ALIGN.RIGHT)


def title(slide, text):
    textbox(slide, text, 0.72, 0.30, 7.0, 0.45, 28, DARK, True)
    bar = slide.shapes.add_shape(1, Inches(0.72), Inches(0.88), Inches(0.82), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def add_image(slide, filename, x, y, w, h):
    path = BASE / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def slide1_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 1)
    textbox(s, "基于时空图神经网络的\n多行为序列推荐模型", 0.95, 1.30, 7.8, 1.20, 32, DARK, True)
    textbox(s, "华南师范大学本科毕业论文答辩", 1.00, 3.00, 5.2, 0.30, 18, GREEN, True)
    textbox(s, "研究主线：多行为数据 → 时空图模型 → 软件落地 → 实验验证", 1.00, 3.52, 7.0, 0.32, 16, MUTED)
    add_image(s, "fig01_stgnn_architecture.png", 8.15, 1.25, 4.35, 3.10)
    panel(s, "研究对象", "围绕点击、收藏、加购、购买等多行为日志，统一建模结构依赖、时间动态和行为强度，并封装为可交互桌面应用。", 0.98, 5.20, 11.0, 0.95, LIGHT, 15, 14)


def slide2_toc():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 2)
    title(s, "目录")
    panel(s, "01 研究背景与技术路线", "问题来源、研究目标、总体路径", 0.88, 1.35, 3.65, 1.10, LIGHT)
    panel(s, "02 理论基础与数据处理", "推荐任务、关键技术、时空图构建", 4.84, 1.35, 3.65, 1.10, WHITE)
    panel(s, "03 模型设计", "空间建模、时间建模、时空融合、训练策略", 8.80, 1.35, 3.65, 1.10, LIGHT)
    panel(s, "04 系统实现", "推荐引擎、桌面交互、质量检验与发布", 0.88, 2.78, 3.65, 1.10, WHITE)
    panel(s, "05 实验结果", "主实验、消融实验、参数敏感性与案例分析", 4.84, 2.78, 3.65, 1.10, LIGHT)
    panel(s, "06 总结与致谢", "研究贡献、后续工作与致谢", 8.80, 2.78, 3.65, 1.10, WHITE)
    add_image(s, "系统总览6.png", 1.20, 4.35, 10.8, 1.75)


def slide3_background():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 3)
    title(s, "研究背景与技术路线")
    panel(s, "研究背景", "电商推荐已从静态展示转向智能分发。用户行为包含点击、收藏、加购、购买等多个阶段，传统单一交互建模难以完整反映兴趣形成过程。", 0.82, 1.42, 5.45, 1.35, LIGHT)
    panel(s, "研究问题", "序列模型擅长时间顺序建模但弱于高阶关系，图模型擅长结构传播但弱于时间动态，多行为语义和时空联合建模仍不充分。", 0.82, 3.00, 5.45, 1.35, WHITE)
    panel(s, "技术路线", "数据清洗与字段统一 → 时空异构图构建 → 时空图模型训练 → 实验分析 → 桌面应用封装。", 0.82, 4.58, 5.45, 1.10, LIGHT)
    add_image(s, "fig03_behavior_distribution.png", 6.85, 1.55, 5.25, 2.55)
    textbox(s, "多行为数据分布", 8.15, 4.20, 2.7, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "研究切入点", "将用户、商品、行为类型和时间信息统一到同一建模框架中，使推荐系统同时理解关系结构与兴趣变化。", 6.65, 4.65, 5.65, 1.05, WHITE)


def slide4_theory_data():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 4)
    title(s, "理论基础与数据处理")
    panel(s, "Top-K任务", "R_u^K = TopK_{i∈I} f(u,i)\n推荐的目标是排序，从候选集合中选出前K个结果。", 0.82, 1.42, 3.60, 1.20, LIGHT)
    panel(s, "关键技术", "序列模型负责时间依赖，图模型负责结构关系，多行为建模负责行为语义强弱表达。", 4.62, 1.42, 3.60, 1.20, WHITE)
    panel(s, "统一字段", "e=(u,i,b,t)\n以用户、商品、行为类型和时间戳四元组统一表示原始行为日志。", 8.42, 1.42, 3.60, 1.20, LIGHT)
    add_image(s, "fig04_user_activity_distribution.png", 0.92, 3.10, 5.15, 2.00)
    add_image(s, "fig05_item_popularity_distribution.png", 6.42, 3.10, 5.15, 2.00)
    textbox(s, "用户活跃度分布", 2.05, 5.15, 2.1, 0.18, 12, GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "商品流行度分布", 7.55, 5.15, 2.1, 0.18, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "数据特征", "数据具有明显长尾特征，点击远多于购买，因此后续模型必须同时处理稀疏性、热门偏置和高低意图行为不均衡问题。", 0.82, 5.52, 11.66, 0.95, WHITE)


def slide5_graph():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 5)
    title(s, "时空异构图构建")
    panel(s, "图定义", "G=(V,E,φ,ψ)\n用户、商品及其关系被组织成带类型的时空异构图，而不是静态同质图。", 0.82, 1.42, 3.95, 1.20, LIGHT)
    panel(s, "时间窗口", "W_k=[t_k, t_k+Δ)\n通过时间窗口切分连续时间轴，降低长跨度噪声对当前兴趣的干扰。", 4.98, 1.42, 3.55, 1.20, WHITE)
    panel(s, "行为强度", "ω_(u,i,t)=w_b·exp(-λΔt)\n同时量化行为强弱与时间远近，用于边权和样本权重构造。", 8.76, 1.42, 3.55, 1.20, LIGHT)
    add_image(s, "fig06_behavior_transition_heatmap.png", 1.10, 3.00, 5.25, 2.35)
    textbox(s, "多行为转移矩阵热力图", 2.20, 5.45, 3.0, 0.18, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "图构建结果", "通过统一字段、清洗规整、时间窗口、动态图快照和行为转移统计，原始日志被整理为可直接进入时空图模型训练的结构化输入。", 6.70, 3.12, 5.55, 1.55, WHITE)


def slide6_model():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 6)
    title(s, "模型设计")
    add_image(s, "fig01_stgnn_architecture.png", 0.90, 1.35, 5.75, 3.45)
    panel(s, "空间建模", "通过用户—商品图和商品转移图进行双图传播，学习用户与商品之间的高阶协同关系。", 6.95, 1.42, 5.15, 1.15, LIGHT)
    panel(s, "时间建模", "序列输入由商品表示、行为嵌入和位置编码组成，并结合时间间隔编码与GRU刻画兴趣变化。", 6.95, 2.82, 5.15, 1.15, WHITE)
    panel(s, "时空融合", "g_u=σ(W_g[h_u^sp||h_u^tm]+b_g)\nz_u=g_u⊙h_u^sp+(1-g_u)⊙h_u^tm\n门控决定当前更依赖结构信号还是时间信号。", 6.95, 4.22, 5.15, 1.40, LIGHT)
    textbox(s, "结构图中的空间分支、时间分支与融合输出", 1.55, 5.00, 4.4, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)


def slide7_training():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 7)
    title(s, "预测层与训练策略")
    panel(s, "预测层", "y_(u,i)=z_u^T h_i^sp\n采用内积形式进行候选打分，表达简洁、计算效率高，适合实时排序。", 0.82, 1.42, 3.95, 1.30, LIGHT)
    panel(s, "BPR损失", "L_bpr=-logσ(y_(u,i+)-y_(u,i-))\n排序学习目标要求正样本得分高于负样本。", 4.98, 1.42, 3.95, 1.30, WHITE)
    panel(s, "在线更新", "z_u^new=ηz_u'+(1-η)z_u^old\n新增行为后只做局部状态刷新，不触发全量重训。", 9.14, 1.42, 3.35, 1.30, LIGHT)
    bullets(s, [
        "训练流程：构图 → 序列打包 → 前向传播 → BPR损失 → 参数更新。",
        "负采样采用随机负采样与难负样本补充相结合的方式。",
        "复杂度主要受图边规模和序列长度影响，训练中通过浅层传播和稀疏矩阵乘法控制开销。",
        "验证集以 NDCG 作为早停依据，用于平衡拟合能力与泛化稳定性。"
    ], 0.95, 3.20, 11.2, 2.60, 17)


def slide8_system():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 8)
    title(s, "系统实现与软件封装")
    panel(s, "系统架构", "系统由数据层、模型层、服务层和界面层组成，整体遵循“界面触发 → 服务调度 → 引擎执行 → 结果回传”的调用链。", 0.82, 1.42, 5.25, 1.25, LIGHT)
    panel(s, "推荐引擎", "统一处理训练、推理、增量更新和历史淘汰。离线路径负责初始训练，在线路径负责新增行为后的局部刷新。", 6.30, 1.42, 6.00, 1.25, WHITE)
    add_image(s, "双工作区_Demo.png", 0.92, 3.00, 5.85, 2.60)
    add_image(s, "推荐刷新-后.png", 6.92, 3.00, 5.35, 2.60)
    textbox(s, "双工作区界面", 2.75, 5.72, 2.0, 0.18, 12, GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "推荐刷新结果", 8.75, 5.72, 2.0, 0.18, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "系统特点", "系统同时支持 Demo 演示环境与 Custom 自定义数据库环境，在统一操作逻辑下完成推荐、交互更新和结果展示。", 0.82, 5.98, 11.66, 0.75, WHITE)


def slide9_experiment():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 9)
    title(s, "实验结果与案例分析")
    add_image(s, "fig08_main_benchmark_metrics.png", 0.88, 1.38, 5.45, 2.35)
    add_image(s, "fig09_ablation_metrics.png", 6.88, 1.38, 5.15, 2.35)
    textbox(s, "主实验结果", 2.55, 3.80, 1.8, 0.18, 12, GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "消融实验结果", 8.55, 3.80, 1.8, 0.18, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "实验结论", "时空联合建模在多数设置下表现更平衡。空间、时间和行为强度三个模块分别承担结构关系、动态变化和意图层次建模，联合使用时效果最稳定。", 0.82, 4.25, 5.45, 1.30, LIGHT)
    panel(s, "案例分析", "结合批量评分、专项诊断和用户行为可视化，可以观察推荐刷新、商品诊断和失败样本归因，从而让模型结果具备更强可解释性。", 6.48, 4.25, 5.80, 1.30, WHITE)


def slide10_thanks():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s, 10)
    textbox(s, "感谢聆听", 4.10, 2.05, 5.2, 0.80, 34, DARK, True, PP_ALIGN.CENTER)
    textbox(s, "恳请各位老师批评指正", 4.00, 3.00, 5.4, 0.35, 18, GREEN, True, PP_ALIGN.CENTER)
    add_image(s, "注册.png", 1.10, 4.05, 3.35, 1.85)
    add_image(s, "Inspector批量评分.png", 4.98, 4.05, 3.35, 1.85)
    add_image(s, "用户购买记录可视化.png", 8.86, 4.05, 3.05, 1.85)


for fn in [
    slide1_cover,
    slide2_toc,
    slide3_background,
    slide4_theory_data,
    slide5_graph,
    slide6_model,
    slide7_training,
    slide8_system,
    slide9_experiment,
    slide10_thanks,
]:
    fn()

try:
    prs.save(PRIMARY_OUT)
    print(PRIMARY_OUT)
except PermissionError:
    prs.save(FALLBACK_OUT)
    print(FALLBACK_OUT)
