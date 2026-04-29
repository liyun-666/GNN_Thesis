# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(r"D:\GNN_Thesis")
PRIMARY_OUT = BASE / "基于时空图神经网络的多行为序列推荐模型答辩ppt.pptx"
FALLBACK_OUT = BASE / "基于时空图神经网络的多行为序列推荐模型答辩ppt_v3.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(248, 250, 252)
GREEN = RGBColor(0, 112, 84)
DARK = RGBColor(26, 43, 52)
MUTED = RGBColor(95, 108, 122)
LIGHT = RGBColor(232, 244, 240)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(208, 225, 220)
GOLD = RGBColor(196, 151, 54)


def textbox(slide, text, x, y, w, h, size=18, color=DARK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def bullets(slide, items, x, y, w, h, size=16, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def panel(slide, title, body, x, y, w, h, fill=WHITE, title_size=15, body_size=13):
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    textbox(slide, title, x + 0.14, y + 0.08, w - 0.28, 0.28, title_size, GREEN, True)
    textbox(slide, body, x + 0.14, y + 0.38, w - 0.28, h - 0.44, body_size, DARK, False)


def bg(slide, index):
    base = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    base.fill.solid()
    base.fill.fore_color.rgb = BG
    base.line.fill.background()

    top = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.14))
    top.fill.solid()
    top.fill.fore_color.rgb = GREEN
    top.line.fill.background()

    left = slide.shapes.add_shape(1, 0, 0, Inches(0.18), prs.slide_height)
    left.fill.solid()
    left.fill.fore_color.rgb = GREEN
    left.line.fill.background()

    deco = slide.shapes.add_shape(9, Inches(11.45), Inches(0.34), Inches(1.45), Inches(1.45))
    deco.fill.solid()
    deco.fill.fore_color.rgb = LIGHT
    deco.line.fill.background()

    textbox(slide, f"{index}/9", 12.05, 6.95, 0.75, 0.22, 10, MUTED, False, PP_ALIGN.RIGHT)


def title(slide, text):
    textbox(slide, text, 0.72, 0.30, 6.4, 0.45, 28, DARK, True)
    bar = slide.shapes.add_shape(1, Inches(0.72), Inches(0.88), Inches(0.82), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def add_image(slide, name, x, y, w, h):
    path = BASE / name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def slide1():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 1)
    title(s, "研究背景与技术路线")
    textbox(s, "多行为推荐的问题来源、研究现状与整体实现路径", 0.74, 1.02, 5.8, 0.28, 15, MUTED)
    panel(s, "研究背景", "电商推荐已从静态展示转向智能分发。用户会产生点击、收藏、加购、购买等多种行为，单一交互信号已难以完整反映兴趣形成过程。", 0.82, 1.42, 5.45, 1.30, LIGHT)
    panel(s, "现有不足", "序列模型擅长顺序建模但弱于高阶关系，图模型擅长结构传播但弱于时间动态，多行为语义差异与时空联合建模仍不充分。", 0.82, 2.92, 5.45, 1.30)
    panel(s, "研究目标", "构建基于时空图神经网络的多行为序列推荐模型，并完成实验验证与桌面系统落地，实现可更新、可解释、可演示的推荐流程。", 0.82, 4.42, 5.45, 1.30)
    panel(s, "技术路线", "数据清洗与字段统一 → 时空异构图构建 → 时空图模型训练 → 主实验/消融/案例分析 → 桌面应用封装。", 0.82, 5.92, 5.45, 1.00, LIGHT)
    frame = s.shapes.add_shape(5, Inches(6.65), Inches(1.42), Inches(5.95), Inches(3.15))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    add_image(s, "fig03_behavior_distribution.png", 6.82, 1.58, 5.62, 2.78)
    textbox(s, "图示：多行为数据分布", 6.82, 4.74, 5.0, 0.22, 13, GREEN, True)
    panel(s, "研究切入点", "研究对象由单一交互扩展为多行为序列，建模重点由静态匹配扩展为结构关系与时间动态的联合学习，并进一步落实到可交互系统实现。", 6.65, 5.52, 5.95, 1.40)


def slide2():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 2)
    title(s, "理论基础与关键技术")
    textbox(s, "推荐任务、评价指标、序列模型、图模型与多行为建模的理论基础", 0.74, 1.02, 7.5, 0.28, 15, MUTED)
    panel(s, "Top-K任务", "R_u^K = TopK_{i∈I} f(u,i)\n排序目标是从候选集合中选出最值得展示的前K个商品，而不是对单个商品做孤立判断。", 0.82, 1.42, 3.95, 1.35, LIGHT, 15, 13)
    panel(s, "评价指标", "HR关注是否命中，NDCG关注命中位置，Recall关注覆盖程度。三者共同衡量“推中了没有”和“排得好不好”。", 4.98, 1.42, 3.95, 1.35, WHITE, 15, 13)
    panel(s, "序列模型", "Attention(Q,K,V)=softmax(QK^T/√d)V\nGRU与Transformer分别代表递推式和注意力式时间建模思路，用于捕捉行为顺序依赖。", 9.14, 1.42, 3.35, 1.35, LIGHT, 15, 13)
    panel(s, "图神经网络", "H^(l+1)=σ(ÂH^(l)W^(l))\n节点表示由自身和邻域共同决定，推荐任务中的价值在于学习用户与商品之间的高阶协同关系。", 0.82, 3.10, 5.10, 1.38, WHITE, 15, 13)
    panel(s, "多行为强度", "ω_(u,i,t)=w_b·exp(-λΔt)\n把行为语义强弱与时间新鲜性同时转化为可学习信号。", 6.15, 3.10, 3.45, 1.38, LIGHT, 15, 13)
    panel(s, "排序学习", "L_BPR = -logσ(y_(u,i+) - y_(u,i-))\n模型优化的是正负样本的相对顺序，这与Top-K推荐的目标一致。", 9.82, 3.10, 2.67, 1.38, WHITE, 15, 12)
    panel(s, "理论支撑关系", "Top-K定义明确了推荐目标，评价指标刻画了排序效果，序列模型提供时间建模基础，图模型提供结构传播基础，多行为建模补足了行为语义与意图差异表达。", 0.82, 4.85, 11.66, 1.48, LIGHT, 15, 14)


def slide3():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 3)
    title(s, "数据处理与字段规范")
    textbox(s, "数据来源、字段定义、清洗流程与长尾分布分析", 0.74, 1.02, 6.2, 0.28, 15, MUTED)
    panel(s, "数据来源", "主数据来自电商用户行为日志，保留点击、收藏、加购、购买四类行为，而不是只保留最终购买结果。", 0.82, 1.42, 3.75, 1.25, LIGHT)
    panel(s, "统一字段", "e=(u,i,b,t)\nu表示用户，i表示商品，b表示行为类型，t表示Unix时间戳。四元组形式是图建模和序列建模的统一输入接口。", 4.80, 1.42, 3.75, 1.25, WHITE)
    panel(s, "清洗目标", "处理缺失值、非法编码、异常时间戳、重复记录和顺序错乱问题，保证数据具备完整性、一致性和可解释性。", 8.78, 1.42, 3.75, 1.25, LIGHT)
    add_image(s, "fig04_user_activity_distribution.png", 0.92, 3.08, 5.15, 2.05)
    add_image(s, "fig05_item_popularity_distribution.png", 6.42, 3.08, 5.15, 2.05)
    textbox(s, "图2 用户活跃度分布", 1.95, 5.18, 2.6, 0.22, 12, GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "图3 商品流行度分布", 7.45, 5.18, 2.6, 0.22, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "图表含义", "两张分布图共同说明数据具有明显长尾特征：少量高活跃用户和头部商品贡献大量交互，大量用户与商品处于长尾区域，因此模型必须处理稀疏性与热门偏置。", 0.82, 5.55, 11.66, 1.18, WHITE, 15, 13)


def slide4():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 4)
    title(s, "时空异构图构建与行为建模")
    textbox(s, "图定义、时间窗口、动态图快照与行为转移机制", 0.74, 1.02, 6.8, 0.28, 15, MUTED)
    panel(s, "图定义", "G=(V,E,φ,ψ)\nV为节点集合，E为边集合，φ和ψ分别表示节点类型映射与边类型映射，说明本文采用的是带类型的时空异构图。", 0.82, 1.42, 4.05, 1.40, LIGHT)
    panel(s, "时间窗口", "W_k=[t_k, t_k+Δ)\n时间窗口把连续时间切分为局部阶段，用于降低长跨度噪声，让模型在不同阶段学习相对稳定的兴趣状态。", 5.08, 1.42, 3.55, 1.40, WHITE)
    panel(s, "动态图快照", "{G^(1), G^(2), …, G^(T)}\n快照序列表达图随时间演化的状态，为新增行为后的局部更新提供了形式化基础。", 8.86, 1.42, 3.62, 1.40, LIGHT)
    add_image(s, "fig06_behavior_transition_heatmap.png", 0.92, 3.12, 5.55, 2.45)
    textbox(s, "图5 多行为转移矩阵热力图", 1.85, 5.62, 3.8, 0.22, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "行为强度函数", "ω_(u,i,t)=w_b·exp(-λΔt)\n它把“行为强弱”和“时间远近”同时转化为可学习信号，用于边权、样本权重和融合输入。", 6.85, 3.10, 5.60, 1.12, WHITE)
    panel(s, "转移概率", "P(b_j|b_i)=N(b_i→b_j)/Σ_k N(b_i→b_k)\n用于量化行为迁移方向，帮助模型关注高频且高价值的意图增强路径。", 6.85, 4.38, 5.60, 1.12, LIGHT)
    panel(s, "关键建模结果", "通过统一字段、清洗规整、图快照定义、行为强度函数和转移概率统计，原始日志被整理成可直接进入时空图模型训练的结构化输入。", 6.85, 5.66, 5.60, 1.05, WHITE)


def slide5():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 5)
    title(s, "问题定义与空间依赖建模")
    textbox(s, "时空推荐任务形式化、异构关系编码与图传播机制", 0.74, 1.02, 6.8, 0.28, 15, MUTED)
    panel(s, "任务形式化", "y_(u,i)(t)=f_Θ(u,i,H_u(<t),G(t))\nR_u^K(t)=TopK_i y_(u,i)(t)\n打分函数同时依赖用户历史序列和时刻t的图状态，说明推荐是时间条件下的动态图排序。", 0.82, 1.42, 5.05, 1.55, LIGHT)
    panel(s, "边权定义", "ω_(u,i)=w_b·exp(-λΔt)\n在用户—商品边上显式保留行为强度和时间衰减，使购买边、近期边在传播中贡献更大。", 6.10, 1.42, 3.90, 1.55, WHITE)
    panel(s, "双图传播", "H_u^(l+1)=σ(Â_ui H_i^(l)W_u^(l)+H_u^(l)W_uu^(l))\nH_i^(l+1)=σ(Â_ii H_i^(l)W_ii^(l)+Â_ui^T H_u^(l)W_iu^(l))\n用户从交互商品邻域聚合信息，商品同时从迁移邻域和用户邻域吸收信号。", 0.82, 3.25, 6.35, 1.75, WHITE)
    panel(s, "层间融合", "h_u^sp=Σ α_l H_u^(l),  h_i^sp=Σ α_l H_i^(l)\n采用浅层传播与层间加权汇总，兼顾高阶关系表达和节点区分能力，降低过平滑风险。", 7.42, 3.25, 5.06, 1.75, LIGHT)
    add_image(s, "fig01_stgnn_architecture.png", 0.98, 5.30, 4.65, 1.20)
    textbox(s, "模型结构图中的空间分支", 1.80, 6.52, 2.9, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "空间模块输出", "空间模块最终输出用户空间表示 h_u^sp 和商品空间表示 h_i^sp。它们提供候选表达基础，回答“哪些商品在结构上与用户相关”。", 6.10, 5.20, 6.38, 1.35, WHITE)


def slide6():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 6)
    title(s, "时间动态建模与时空融合")
    textbox(s, "顺序编码、时间间隔建模、长短期兴趣融合与动态状态更新", 0.74, 1.02, 7.4, 0.28, 15, MUTED)
    panel(s, "顺序输入", "x_l=h_i^sp+e_b+p_l\n每一步输入由商品表示、行为嵌入和位置编码组成，使同一商品在不同类型行为下具有不同语义。", 0.82, 1.42, 3.90, 1.40, LIGHT)
    panel(s, "时间间隔编码", "g_l=φ(log(1+Δt_l)),  x_l=x_l+g_l\n对时间差取对数后再映射，既压缩极端长间隔，又让模型感知行为远近。", 4.95, 1.42, 3.90, 1.40, WHITE)
    panel(s, "长短期兴趣", "h_u^tm=W_s h_u^short + W_l h_u^long\n短期兴趣来自序列末端状态，长期兴趣来自序列隐藏状态统计，两者共同构成时间表示。", 9.08, 1.42, 3.40, 1.40, LIGHT)
    panel(s, "融合门控", "g_u=σ(W_g[h_u^sp||h_u^tm]+b_g)\nz_u=g_u⊙h_u^sp + (1-g_u)⊙h_u^tm\n门控决定当前更依赖结构信号还是时间信号。", 0.82, 3.15, 5.35, 1.65, WHITE)
    panel(s, "状态更新", "z_u^new=ηz_u' + (1-η)z_u^old\n新增行为到来后只更新受影响用户的局部状态，以降低波动并支持在线刷新。", 6.40, 3.15, 3.00, 1.65, LIGHT)
    panel(s, "语义对齐", "c_u=Σ π_b e_b,  z_u=z_u+W_c c_u\n近期行为结构被压缩成语义上下文向量，对融合后的用户状态进行校正。", 9.63, 3.15, 2.85, 1.65, WHITE)
    panel(s, "融合结果", "z_u 是最终用户动态兴趣表示。它同时包含长期结构偏好、短期行为变化和行为强弱信息，是后续推荐打分的直接输入。", 0.82, 5.15, 11.66, 1.18, LIGHT)


def slide7():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 7)
    title(s, "预测层与训练策略")
    textbox(s, "评分函数、BPR损失、负采样机制与复杂度控制", 0.74, 1.02, 6.8, 0.28, 15, MUTED)
    panel(s, "预测层", "y_(u,i)=z_u^T h_i^sp\n采用内积形式进行候选打分，原因是表达简洁、计算效率高，适合在线系统中的实时排序任务。", 0.82, 1.42, 4.10, 1.45, LIGHT)
    panel(s, "BPR损失", "L_bpr=-logσ(y_(u,i+)-y_(u,i-))\nL=-w_b+ logσ(y_(u,i+)-y_(u,i-)) + β||Θ||_2^2\n高意图行为拥有更大权重，排序优化聚焦正负样本相对次序。", 5.15, 1.42, 4.05, 1.45, WHITE)
    panel(s, "负采样", "训练时结合随机负采样与难负样本补充。随机采样保证覆盖面，难负样本提高判别强度，同时按时间约束避免未来信息泄漏。", 9.43, 1.42, 3.05, 1.45, LIGHT)
    bullets(s, [
        "训练流程：构图 → 序列打包 → 前向传播 → BPR损失 → 反向更新。",
        "在线流程：新增行为写入后执行局部更新，不触发每次全量重训。",
        "复杂度主要由图边规模 |E| 和序列长度 L_s 决定，训练中通过浅层传播、稀疏矩阵乘法和批处理控制开销。",
        "验证集以 NDCG 作为早停依据，用于平衡拟合能力与泛化稳定性。"
    ], 0.96, 3.30, 11.2, 2.55, 17)
    panel(s, "训练目标对应关系", "预测层回答“用户对哪个商品更可能感兴趣”，BPR与负采样回答“怎样把真正相关的商品排到更前面”，复杂度控制则保证模型能够真正运行在交互式应用环境中。", 0.82, 5.95, 11.66, 0.78, WHITE, 15, 13)


def slide8():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 8)
    title(s, "系统架构与推荐引擎")
    textbox(s, "分层架构、模块调用关系、离线训练与在线更新机制", 0.74, 1.02, 6.8, 0.28, 15, MUTED)
    panel(s, "系统分层", "系统由数据层、模型层、服务层和界面层组成。数据层负责输入约束与存储，模型层承担推荐计算，服务层负责调度训练、推理和更新，界面层负责参数输入和结果反馈。", 0.82, 1.42, 5.15, 1.45, LIGHT)
    panel(s, "模块调用", "整体调用链遵循“界面触发 → 服务调度 → 引擎执行 → 结果回传”。这样设计可以减少界面与底层模型的直接耦合，也便于扩展训练、诊断和导出功能。", 6.20, 1.42, 6.28, 1.45, WHITE)
    add_image(s, "双工作区_Demo.png", 0.92, 3.18, 5.85, 2.70)
    add_image(s, "推荐刷新-后.png", 6.92, 3.18, 5.35, 2.70)
    textbox(s, "双工作区界面", 2.75, 5.95, 2.0, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "推荐刷新结果", 8.75, 5.95, 2.0, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "引擎实现要点", "推荐引擎统一处理训练、推理、增量更新和历史淘汰。离线路径负责初始训练与周期性重训，在线路径负责用户新增行为后的局部状态刷新，两条路径共享相同的打分和排序逻辑。", 0.82, 6.12, 11.66, 0.70, WHITE, 15, 13)


def slide9():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 9)
    title(s, "桌面交互、诊断与发布")
    textbox(s, "登录管理、批量评分、用户可视化与程序打包", 0.74, 1.02, 6.0, 0.28, 15, MUTED)
    add_image(s, "注册.png", 0.88, 1.40, 3.45, 1.90)
    add_image(s, "Inspector批量评分.png", 4.72, 1.40, 3.95, 1.90)
    add_image(s, "用户购买记录可视化.png", 8.98, 1.40, 3.30, 1.90)
    textbox(s, "注册登录", 1.95, 3.35, 1.2, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "Inspector批量评分", 5.95, 3.35, 1.7, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "用户行为可视化", 9.95, 3.35, 1.7, 0.20, 12, GREEN, True, PP_ALIGN.CENTER)
    panel(s, "交互功能", "桌面应用把推荐、行为提交、数据库接入、诊断与导出整合到同一界面环境中。Demo工作区用于稳定演示，Custom工作区支持外部数据库加载与训练。", 0.82, 3.72, 3.75, 1.38, LIGHT)
    panel(s, "质量检验", "Inspector支持批量抽样评分与专项诊断，可输出cases、pass、avgQ以及逐条样本的分数变化、排名变化和诊断消息，用于分析推荐是否稳定合理。", 4.78, 3.72, 3.75, 1.38, WHITE)
    panel(s, "发布方式", "系统基于Python实现，可打包为单文件可执行程序和安装包。模型文件、数据库文件与配置资源统一纳入产物中，使论文中的方法能够真正以桌面应用方式运行。", 8.74, 3.72, 3.75, 1.38, LIGHT)
    panel(s, "系统实现结果", "第五章的重点不是重新解释模型，而是说明模型如何被组织进可运行的软件框架：数据能接入、推荐能刷新、结果能诊断、系统能发布。这使论文工作从算法验证进一步走向了应用落地。", 0.82, 5.40, 11.66, 1.18, WHITE, 15, 13)


for fn in [slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8, slide9]:
    fn()

try:
    prs.save(PRIMARY_OUT)
    print(PRIMARY_OUT)
except PermissionError:
    prs.save(FALLBACK_OUT)
    print(FALLBACK_OUT)
